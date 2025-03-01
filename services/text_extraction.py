import os
import pymupdf
from PIL import Image
from io import BytesIO
from docx import Document  # For .docx files
from pptx import Presentation  # For .pptx files
import requests
from bs4 import BeautifulSoup
import re
from youtube_transcript_api import YouTubeTranscriptApi
import base64
from google import genai

# Load the Gemini API key from the environment variable
api_key = os.getenv("GEMINI_API_KEY")

# Ensure the API key is found
if not api_key:
    raise ValueError("API key for Gemini not found. Set the GEMINI_API_KEY environment variable.")

# Replace Tesseract with Gemini for OCR
def extract_text_from_image_bytes(image_bytes: bytes ,code) -> str:
    # Initialize Gemini client with the loaded API key
    client = genai.Client(api_key=api_key)

    # Convert the image to base64 encoding
    img_base64 = base64.b64encode(image_bytes).decode("utf-8")
    
    # Send the image to Gemini for text generation
    response = client.models.generate_content(
        model="gemini-2.0-flash",  # Or the model you intend to use
        contents=[img_base64, "Extract text from this image"]
    )
    
    # Check if the response contains the expected text
    if response.text:
        return response.text
    else:
        return "No text extracted from the image."

# Extract text from PDFs
def extract_text_from_pdf_bytes(pdf_bytes: bytes,code) -> str:
    doc = pymupdf.open(stream=pdf_bytes, filetype="pdf")
    return "\n".join(page.get_text() for page in doc)

# Extract text from Word documents (.docx)
def extract_text_from_docx_bytes(docx_bytes: bytes,code) -> str:
    doc = Document(BytesIO(docx_bytes))
    return "\n".join(para.text for para in doc.paragraphs)

# Extract text from PowerPoint files (.pptx)
def extract_text_from_pptx_bytes(pptx_bytes: bytes,code) -> str:
    presentation = Presentation(BytesIO(pptx_bytes))
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# Extract text from plain text (.txt)
def extract_text_from_txt_bytes(txt_bytes: bytes,code) -> str:
    return txt_bytes.decode("utf-8")

# Extract text from a webpage (web scraping)
def extract_text_from_url(url: str) -> str:
    response = requests.get(url)
    if response.status_code == 200:
        soup = BeautifulSoup(response.content, "html.parser")
        return soup.get_text()
    return "Failed to extract text"

# Extract YouTube transcript
def extract_youtube_transcript(video_url: str,code) -> str:
    # Regular expression to extract video ID from various types of YouTube links
    video_id_match = re.search(r"v=([a-zA-Z0-9_-]+)", video_url)
    
    if video_id_match:
        video_id = video_id_match.group(1)
    else:
        return "Invalid YouTube URL or video ID not found"
    
    try:
        # Attempt to fetch the transcript in English first, then fallback to French
        transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=['en', 'fr'])
        return "\n".join([entry["text"] for entry in transcript])
    except Exception as e:
        return f"Transcript not available. Error: {str(e)}"



def get_book_data(code, folder_path="content_materials/"):
    """
    Finds the most recent text file in the folder whose name starts with the given code followed by '_'
    and returns its content.
    
    :param code: The code to match at the beginning of the filename.
    :param folder_path: The path to the folder containing text files.
    :return: The content of the matched text file or an error message.
    """
    if not os.path.exists(folder_path):
        return "Folder does not exist."

    # List all text files that start with 'code_'
    matching_files = [
        f for f in os.listdir(folder_path) 
        if f.startswith(f"{code}_") and f.endswith(".txt")
    ]

    if not matching_files:
        return f"No matching files found for code: {code}"

    # Sort by creation time (newest first)
    matching_files.sort(key=lambda x: os.path.getctime(os.path.join(folder_path, x)), reverse=True)

    latest_file = os.path.join(folder_path, matching_files[0])

    # Read and return the content of the latest matching file
    with open(latest_file, "r", encoding="utf-8") as file:
        return file.read()

